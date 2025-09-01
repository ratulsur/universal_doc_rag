# ingestor/image_extractor.py
from __future__ import annotations
import sys
import io
import shutil
from pathlib import Path
from typing import List, Optional

from PIL import Image
import pdfplumber
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException


class ImageExtractor:
    SUPPORTED = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}

    def __init__(self, out_dir: str | Path = "data/extracted_images") -> None:
        self.log = CustomLogger.get_logger(__name__)
        self.out_dir = Path(out_dir)
        self.out_dir.mkdir(parents=True, exist_ok=True)

    def extract(self, file_path: str | Path, prefix: Optional[str] = None) -> List[Path]:
        """
        Returns list of image file paths saved under out_dir.
        """
        try:
            path = Path(file_path)
            ext = path.suffix.lower()
            tag = prefix or path.stem

            if ext not in self.SUPPORTED:
                self.log.warning("Unsupported for image extraction: %s", ext)
                return []

            if ext == ".pdf":
                return self._from_pdf(path, tag)
            if ext == ".docx":
                return self._from_docx(path, tag)
            if ext == ".pptx":
                return self._from_pptx(path, tag)
            # regular image: copy
            return [self._copy_image(path, tag)]
        except Exception as e:
            self.log.error("Failed to extract images: %s", e)
            raise DocumentPortalException("Image extraction error", sys) from e

    # ---------- Implementations ----------
    def _from_pdf(self, path: Path, tag: str) -> List[Path]:
        saved: List[Path] = []
        
        try:
            with pdfplumber.open(str(path)) as pdf:
                for pnum, page in enumerate(pdf.pages, start=1):
                    for img in page.images:
                        try:
                            # cuz pdfplumber exposes image box; we need raw stream -> fallback to fitz for bytes
                            
                            pass
                        except Exception:
                            continue
        except Exception as e:
            self.log.debug("pdfplumber image pass failed: %s (will use fitz)", e)

        
        try:
            doc = fitz.open(str(path))
            for i in range(len(doc)):
                page = doc[i]
                for img_idx, img in enumerate(page.get_images(full=True), start=1):
                    xref = img[0]
                    base = f"{tag}_p{i+1}_{img_idx}.png"
                    out_path = self.out_dir / base
                    pix = fitz.Pixmap(doc, xref)
                    if pix.alpha:  # handle RGBA
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    pix.save(str(out_path))
                    saved.append(out_path)
        except Exception as e:
            self.log.warning("PyMuPDF image extraction failed: %s", e)

        self.log.info("PDF images extracted: %s | file=%s", len(saved), path.name)
        return saved

    def _from_docx(self, path: Path, tag: str) -> List[Path]:
        saved: List[Path] = []
        doc = DocxDocument(str(path))
        
        media_dir = path.parent / f"{path.stem}_media_cache"
        media_dir.mkdir(exist_ok=True)
        try:
            for i, rel in enumerate(doc.part._rels.values(), start=1):  
                if "image" in rel.target_ref:  
                    bin_data = rel.target_part.blob  
                    out_path = self.out_dir / f"{tag}_img_{i}.png"
                    Image.open(io.BytesIO(bin_data)).save(out_path)
                    saved.append(out_path)
        except Exception as e:
            self.log.debug("DOCX rel-scan failed: %s", e)

        self.log.info("DOCX images extracted: %s | file=%s", len(saved), path.name)
        return saved

    def _from_pptx(self, path: Path, tag: str) -> List[Path]:
        saved: List[Path] = []
        prs = Presentation(str(path))
        for sidx, slide in enumerate(prs.slides, start=1):
            for shidx, shape in enumerate(slide.shapes, start=1):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    out_path = self.out_dir / f"{tag}_s{sidx}_{shidx}.{image.ext or 'png'}"
                    with open(out_path, "wb") as f:
                        f.write(image.blob)
                    saved.append(out_path)
        self.log.info("PPTX images extracted: %s | file=%s", len(saved), path.name)
        return saved

    def _copy_image(self, path: Path, tag: str) -> Path:
        out_path = self.out_dir / f"{tag}{path.suffix.lower()}"
        shutil.copy2(path, out_path)
        return out_path
if __name__ == "__main__":
    print("Image extractor module has been loaded successfully")
