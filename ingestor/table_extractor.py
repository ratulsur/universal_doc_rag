# ingestor/table_extractor.py
from __future__ import annotations
import sys
from pathlib import Path
from typing import List, Tuple, Optional
import pandas as pd

from logger.custom_logger import CustomLogger
from exception.custom_exception import DocumentPortalException

# Third-party readers (import lazily where helpful)
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation


class TableExtractor:
    SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".txt", ".md"}

    def __init__(self) -> None:
        self.log = CustomLogger.get_logger(__name__)

    # ---------- Public API given here ----------
    def extract(self, file_path: str | Path) -> List[pd.DataFrame]:
        """Return a list of DataFrames (one per detected table)"""
        try:
            path = Path(file_path)
            ext = path.suffix.lower()

            if ext not in self.SUPPORTED:
                self.log.warning("Unsupported for table extraction: %s", ext)
                return []

            if ext == ".pdf":
                return self._from_pdf(path)
            if ext == ".docx":
                return self._from_docx(path)
            if ext == ".pptx":
                return self._from_pptx(path)
            if ext == ".xlsx":
                return self._from_xlsx(path)
            if ext == ".csv":
                return [pd.read_csv(path)]
            if ext in (".txt", ".md"):
                return self._from_text_like(path)

            return []
        except Exception as e:
            self.log.error("Failed to extract tables: %s", e)
            raise DocumentPortalException("Table extraction error", sys) from e

    # ---------- Implementations ----------#
    def _from_pdf(self, path: Path) -> List[pd.DataFrame]:
        dfs: List[pd.DataFrame] = []
        with pdfplumber.open(str(path)) as pdf:
            for page_idx, page in enumerate(pdf.pages, start=1):
                try:
                    tables = page.extract_tables() or []
                    for t in tables:
                        
                        if t and all(v is not None for v in t[0]):
                            df = pd.DataFrame(t[1:], columns=t[0])
                        else:
                            df = pd.DataFrame(t)
                        dfs.append(df)
                except Exception as e:
                    self.log.warning("PDF table parse error p%s: %s", page_idx, e)
        self.log.info("PDF tables extracted: %s | file=%s", len(dfs), path.name)
        return dfs

    def _from_docx(self, path: Path) -> List[pd.DataFrame]:
        dfs: List[pd.DataFrame] = []
        doc = DocxDocument(str(path))
        for ti, table in enumerate(doc.tables, start=1):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if not rows:
                continue
            if len({len(r) for r in rows}) > 1:
                
                maxw = max(len(r) for r in rows)
                rows = [r + [""] * (maxw - len(r)) for r in rows]
            # header guesses
            header = rows[0] if rows else []
            body = rows[1:] if len(rows) > 1 else []
            try:
                df = pd.DataFrame(body, columns=header if all(h for h in header) else None)
            except Exception:
                df = pd.DataFrame(rows)
            dfs.append(df)
        self.log.info("DOCX tables extracted: %s | file=%s", len(dfs), path.name)
        return dfs

    def _from_pptx(self, path: Path) -> List[pd.DataFrame]:
        dfs: List[pd.DataFrame] = []
        prs = Presentation(str(path))
        for si, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not hasattr(shape, "has_table"):
                    continue
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for r in table.rows:
                        rows.append([c.text_frame.text.strip() if c.text_frame else "" for c in r.cells])
                    if not rows:
                        continue
                    header = rows[0]
                    body = rows[1:] if len(rows) > 1 else []
                    try:
                        df = pd.DataFrame(body, columns=header if all(h for h in header) else None)
                    except Exception:
                        df = pd.DataFrame(rows)
                    dfs.append(df)
        self.log.info("PPTX tables extracted: %s | file=%s", len(dfs), path.name)
        return dfs

    def _from_xlsx(self, path: Path) -> List[pd.DataFrame]:
        try:
            xl = pd.ExcelFile(path)
            dfs = [xl.parse(sheet) for sheet in xl.sheet_names]
            self.log.info("XLSX sheets extracted: %s | file=%s", len(dfs), path.name)
            return dfs
        except Exception:
            # this helps to fall back to a single read if engine issues
            return [pd.read_excel(path)]

    def _from_text_like(self, path: Path) -> List[pd.DataFrame]:
        """
        Heuristic: detect simple pipe/CSV/TSV tables within text/markdown.
        For real markdown grid tables consider adding 'markdown-it-py' + a plugin later.
        """
        lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
        # very normal detector: rows with pipes AND roughly equal splits
        rows = [ln for ln in lines if "|" in ln]
        if not rows:
            return []
        split_rows = [ [c.strip() for c in r.split("|")] for r in rows ]
        width = max(len(r) for r in split_rows)
        split_rows = [ r + [""]*(width-len(r)) for r in split_rows ]
        header = split_rows[0]
        body = split_rows[1:]
        try:
            return [pd.DataFrame(body, columns=header if all(header) else None)]
        except Exception:
            return [pd.DataFrame(split_rows)]

if __name__ == "__main__":
    print("Yay!!Table extractor module loaded successfully")
